# 导入需要的库
from pptx import Presentation
from pptx.util import Inches

# 创建一个新的PPTX文件，并添加一个标题幻灯片
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "中国人教版初中地理八年级下册西北地区"
subtitle.text = "制作人：XXX"

# 添加第一部分：西北地区概述
content_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(content_slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "第一部分：西北地区概述"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "西北地区的概念和范围"
p.level = 1
p = tf.add_paragraph()
p.text = "西北地区的地理位置和边境线"
p.level = 2
p = tf.add_paragraph()
p.text = "西北地区的自然条件：气候、地貌和水文"
p.level = 2

# 添加第二部分：西北地区的资源
slide = prs.slides.add_slide(content_slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "第二部分：西北地区的资源"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "西北地区的矿产资源：煤、石油、天然气、铀等"
p.level = 1
p = tf.add_paragraph()
p.text = "西北地区的农业资源：小麦、玉米、棉花等"
p.level = 2
p = tf.add_paragraph()
p.text = "西北地区的畜牧资源：羊、牛等"
p.level = 2

# 添加第三部分：西北地区的经济发展
slide = prs.slides.add_slide(content_slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "第三部分：西北地区的经济发展"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "西北地区的经济特点和历史"
p.level = 1
p = tf.add_paragraph()
p.text = "西北地区的产业结构和发展现状"
p.level = 2
p = tf.add_paragraph()
p.text = "西北地区的发展前景和挑战"
p.level = 2

# 添加第四部分：西北地区的旅游资源
slide = prs.slides.add_slide(content_slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1]
title.text = "第四部分西北地区的旅游资源"
tf = body.text_frame
p = tf.add_paragraph()
p.text = "西北地区的旅游资源概述"
p.level = 1
p = tf.add_paragraph()
p.text = "西北地区的自然景观：大漠、戈壁、草原等"
p.level = 2
p = tf.add_paragraph()
p.text = "西北地区的文化景观：丝绸之路、古长城、秦始皇兵马俑等"
p.level = 2




prs.save("西北地区.pptx")
